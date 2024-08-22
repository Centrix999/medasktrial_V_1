# views.py
from django.shortcuts import render, redirect
from django.contrib.auth import login, logout

from django.contrib.auth.decorators import login_required
from django.contrib.auth.forms import UserCreationForm, AuthenticationForm
from app.forms import SignUpForm,New_UploadFileForm,New_Prompt
from app.models import Chunk,UploadFile,Prompt
import docx 
import PyPDF2
import re

from django.shortcuts import render, get_object_or_404

#from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback

from django.conf import settings
from operator import attrgetter



from django.contrib import messages
from django.http import JsonResponse

import pyperclip

from concurrent.futures import ThreadPoolExecutor
from langchain.prompts import PromptTemplate




from langchain.text_splitter import RecursiveCharacterTextSplitter

from PyPDF2 import PdfReader

import re
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer


from docx import Document
from django.http import HttpResponse
from django.utils.text import slugify
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import logging
from pptx import Presentation





import os
from pptx.util import Pt
from pptx import Presentation
from pptx.util import Pt
import os
import logging
from django.http import HttpResponse
from django.utils.text import slugify
from django.conf import settings
from django.contrib.auth.decorators import login_required




from pptx import Presentation
from pptx.util import Pt

from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from django.shortcuts import render
from django.http import HttpResponse
from django.contrib.auth.decorators import login_required
from django.conf import settings
import os
import logging
from slugify import slugify
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from django.shortcuts import render
from django.http import HttpResponse
from django.contrib.auth.decorators import login_required
from django.conf import settings
import os
import logging
from slugify import slugify



from django.shortcuts import render, redirect
from django.contrib.auth import login, logout
from django.contrib.auth.forms import AuthenticationForm
from django.utils import timezone
from .models import UserProfile, UserSession

from reportlab.lib.pagesizes import letter
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from reportlab.lib.styles import getSampleStyleSheet


from django.contrib.auth.decorators import login_required
from django.http import HttpResponse
from django.shortcuts import render
from django.utils.text import slugify
from docx import Document
import logging


from django.conf import settings
import openai
import faiss
import numpy as np
import time



import json
from django.http import HttpResponse
from django.shortcuts import render
from django.utils.text import slugify
# Ensure OpenAI API key is set in Django settings
openai.api_key = settings.OPENAI_API_KEY


from sentence_transformers import SentenceTransformer
from django.conf import settings
from sklearn.metrics.pairwise import cosine_similarity



# Initialize OpenAI API
openai.api_key = settings.OPENAI_API_KEY

# Load Sentence Transformer model
model = SentenceTransformer('paraphrase-MiniLM-L6-v2')

def signup_view(request):
    if request.method == 'POST':
        form = SignUpForm(request.POST)
        if form.is_valid():
            user = form.save()
            login(request, user)
            return redirect('redirect_to_dashboard')
    else:
        form = SignUpForm()
    return render(request, 'app/signup.html', {'form': form})






def login_view(request):
    if request.method == 'POST':
        form = AuthenticationForm(data=request.POST)
        if form.is_valid():
            user = form.get_user()
            login(request, user)
            # Record login time in UserSession
            user_profile, created = UserProfile.objects.get_or_create(user=user)
            UserSession.objects.create(user_profile=user_profile, login_time=timezone.now())
            return redirect('redirect_to_dashboard')
    else:
        form = AuthenticationForm()

    form.fields['username'].widget.attrs['placeholder'] = 'Enter your username'
    form.fields['password'].widget.attrs['placeholder'] = 'Enter your password'

    return render(request, 'app/login.html', {'form': form})

def logout_view(request):
    if request.method == 'POST':
        try:
            user_profile = UserProfile.objects.get(user=request.user)
            # Get the latest session for the user profile
            user_session = UserSession.objects.filter(user_profile=user_profile).latest('login_time')
            user_session.logout_time = timezone.now()
            user_session.save()
        except UserProfile.DoesNotExist:
            # Handle the case where UserProfile does not exist
            pass
        except UserSession.DoesNotExist:
            # Handle the case where UserSession does not exist
            pass

        logout(request)
        return redirect('login')


def home(request):
    return render(request, "app/home.html")





@login_required
def redirect_to_dashboard(request):
    if request.user.is_staff:
        return redirect('admin_dashboard')
    else:
        return redirect('user_dashboard')  # Assuming you have a user dashboard view named 'user_dashboard'
    

@login_required
def all_docs(request):
    docs=UploadFile.objects.all()
    context={"files":docs}
    return render(request,'admin/list_docs.html',context)





# new 'edit jul 30

@login_required
def edit_doc(request, doc_id):
    doc = get_object_or_404(UploadFile, pk=doc_id)
    
    upload_form = New_UploadFileForm(instance=doc)
    prompt_form = New_Prompt()
    
    prompts = doc.prompt.all()
    
    if request.method == 'POST':
        # initial_chunk_size = 10000  # Initial large chunk size
        # recursive_chunk_size = 1500  # Final smaller chunk size
        # chunk_overlap = 300
        initial_chunk_size = 4000  # Initial large chunk size
        recursive_chunk_size = 700  # Final smaller chunk size
        chunk_overlap = 300
        
        if 'upload_submit' in request.POST:  
            upload_form = New_UploadFileForm(request.POST, request.FILES, instance=doc)
            if upload_form.is_valid():
                upload = upload_form.save(commit=False)
                upload.uploaded_by = request.user
                upload.save()
                try:

                    with open(upload.file.path, 'rb') as pdf_file:
                        pdf_reader = PdfReader(pdf_file)
                        total_pages = len(pdf_reader.pages)

                        for page_number in range(total_pages):
                            page = pdf_reader.pages[page_number]
                            text = page.extract_text()
                            cleaned_text = text_clean(text)  # Clean the text

                            # Initial chunking
                            initial_chunks = initial_chunking(cleaned_text, initial_chunk_size)

                            for initial_chunk in initial_chunks:
                                # Use RecursiveCharacterTextSplitter
                                splitter = RecursiveCharacterTextSplitter(
                                    chunk_size=recursive_chunk_size,
                                    chunk_overlap=chunk_overlap,
                                    separators=["\n\n", "\n", " "]
                                )
                                final_chunks = splitter.split_text(initial_chunk)

                                for chunk_text in final_chunks:
                                    # Generate embeddings for the chunk
                                    embeddings = get_embeddings(chunk_text)

                                    # Store the chunk with its embeddings
                                    Chunk.objects.create(
                                        document=upload,
                                        content=chunk_text,
                                        embeddings=embeddings,
                                        page_number=page_number + 1
                                    )

                    return redirect('list_doc')
                except Exception as e:
                    print(f"Error processing PDF: {e}")
                    # Optionally, delete the upload record if processing fails
                    upload.delete()
                    # Optionally, add a user-friendly error message
                    messages.error(request, 'Failed to process the PDF. Please try again.')

            else:
                print(upload_form.errors)
                messages.error(request, upload_form.errors)
        
        elif 'prompt_submit' in request.POST:  
            prompt_form = New_Prompt(request.POST)
            if prompt_form.is_valid():
                prompt = prompt_form.save(commit=False)
                prompt.document = doc  
                prompt.save()
                return redirect('edit_doc', doc_id=doc_id) 
            else:
                print(prompt_form.errors)
                messages.error(request, prompt_form.errors)
            
    return render(request, 'admin/edit_doc.html', {'form': upload_form, 
                                                   'prompt_form': prompt_form,
                                                   'doc': doc,
                                                   'prompts': prompts})


# user side


@login_required
def edit_prompt(request, prompt_id):
    prompt = get_object_or_404(Prompt, id=prompt_id)
    if request.method == 'POST':
        form = New_Prompt(request.POST, instance=prompt)
        if form.is_valid():
            form.save()
            return redirect('edit_doc', doc_id=prompt.document.id)  # Redirect back to the document editing page
    else:
        form = New_Prompt(instance=prompt)
    return render(request, 'admin/editing_prompt.html', {'form': form})


@login_required
def delete_prompt(request, prompt_id):
    prompt = get_object_or_404(Prompt, id=prompt_id)
    doc_id = prompt.document.id  # Storing the document ID before deleting the prompt
    if request.method == 'POST':
        prompt.delete()  # Deleting the prompt
        return redirect('edit_doc', doc_id=doc_id)  # Redirecting to the document editing page
    return render(request, 'admin/delete_prompt.html', {'prompt': prompt})

@login_required
def confirm_delete(request,doc_id):
    doc=get_object_or_404(UploadFile,pk=doc_id)
    if request.method == 'POST':
        doc.delete()
        return redirect("list_doc")
    return render(request,'admin/delete_doc_confirm.html',{'doc':doc})



@login_required
def user_dashboard(request):
   


    username = request.user.username
    first_name = request.user.first_name
    last_name = request.user.last_name
    email = request.user.email
    documents = UploadFile.objects.all()
    context={"documents":documents,
             "username":username,
             "first_name":first_name,
             "last_name":last_name,
             "email":email,}
    return render(request, 'user/user_dashboard.html', context)




@login_required
def admin_dashboard(request):
    username = request.user.username
    first_name = request.user.first_name
    last_name = request.user.last_name
    email = request.user.email
    documents = UploadFile.objects.all()
    context={"files":documents,
             "username":username,
             "first_name":first_name,
             "last_name":last_name,
             "email":email,}
    return render(request, 'admin/admin_dashboard.html', context)    



def text_clean(text):
    # Replace "%" with "percentage"
    cleaned_text = text.replace("%", "percentage")
    print('------------after replacing  the  percentage sign-----------')
    print(cleaned_text)
    
    # Remove excess whitespace using regular expressions
    print('------------cleaned_text with out white spaces---')
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
    print(cleaned_text,'\n')
    
    # Lowercase the text
    print('--------------------after lower casing ----------------------')
    lowercased_text = cleaned_text.lower()
    print(lowercased_text)

    # Remove special characters
    # cleaned_text_without_special_chars = re.sub(r'[^A-Za-z0-9\s]', '', lowercased_text)
    cleaned_text_without_special_chars = re.sub(r'[^A-Za-z0-9\s./$₹]', '', lowercased_text)
    print('-----------cleaned_text_without_special_chars')
    print(cleaned_text_without_special_chars,'---------------------cleaned_text_without_special_chars')

    # Define stopwords
    stop_words = set(stopwords.words('english'))
    print(stop_words,'-----------stop_words')
    print('\n')

    # Tokenize the text
    words = word_tokenize(cleaned_text_without_special_chars)
    print(words,'-----------Tokenizewords')
    print('\n')

    # Remove stopwords
    filtered_words = [word for word in words if word not in stop_words]
    print('----------------------------')
    print(filtered_words)
    print('---------------end filterd  words')
    print('\n')

    # Join the filtered words back into a sentence
    cleaned_text_without_stopwords = ' '.join(filtered_words)
    print(cleaned_text_without_stopwords,'---------------------join the text with out stopword')
    print('\n')

    # Initialize the lemmatizer
    lemmatizer = WordNetLemmatizer()
    print('---------',lemmatizer,'---------------------lemmatizer')

    # Lemmatize the text
    lemmatized_text= [lemmatizer.lemmatize(word) for word in word_tokenize(cleaned_text_without_stopwords)]
    print(lemmatized_text,'----------------------lemmatized_text')
    print('\n')
    print(' '.join(lemmatized_text),'-----------\n')

    return ' '.join(lemmatized_text)




import openai
def initial_chunking(text, chunk_size):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]


from sentence_transformers import SentenceTransformer


def admin_new_upload(request):
    # initial_chunk_size = 10000  # Initial large chunk size
    # recursive_chunk_size = 1500  # Final smaller chunk size
    # chunk_overlap = 300
    initial_chunk_size = 4000  # Initial large chunk size
    recursive_chunk_size = 1500  # Final smaller chunk size
    chunk_overlap = 300

    if request.method == 'POST':
        form = New_UploadFileForm(request.POST, request.FILES)
        if form.is_valid():
            upload = form.save(commit=False)
            upload.uploaded_by = request.user
            upload.save()

            try:
                with open(upload.file.path, 'rb') as pdf_file:
                    pdf_reader = PdfReader(pdf_file)
                    total_pages = len(pdf_reader.pages)

                    for page_number in range(total_pages):
                        page = pdf_reader.pages[page_number]
                        text = page.extract_text()
                        cleaned_text = text_clean(text)  # Clean the text

                        # Initial chunking
                        initial_chunks = initial_chunking(cleaned_text, initial_chunk_size)

                        for initial_chunk in initial_chunks:
                            # Use RecursiveCharacterTextSplitter
                            splitter = RecursiveCharacterTextSplitter(
                                chunk_size=recursive_chunk_size,
                                chunk_overlap=chunk_overlap,
                                separators=["\n\n", "\n", " "]
                            )
                            final_chunks = splitter.split_text(initial_chunk)

                            for chunk_text in final_chunks:
                                # Generate embeddings for the chunk
                                embeddings = get_embeddings(chunk_text)

                                # Store the chunk with its embeddings
                                Chunk.objects.create(
                                    document=upload,
                                    content=chunk_text,
                                    embeddings=embeddings,
                                    page_number=page_number + 1
                                )

                return redirect('list_doc')
            except Exception as e:
                # Handle exceptions (e.g., log the error, display a message to the user)
                print(f"Error processing PDF: {e}")
                # Optionally, delete the upload record if processing fails
                upload.delete()
                # Optionally, add a user-friendly error message
                return render(request, 'admin/admin_new_upload.html', {'form': form, 'error': 'Failed to process the PDF. Please try again.'})
    else:
        form = New_UploadFileForm()

    return render(request, 'admin/admin_new_upload.html', {'form': form})







def get_embeddings(text):
    start_time = time.time()
    
    # Generate embeddings
    embeddings = model.encode([text], convert_to_tensor=True)
    
    # Convert to numpy array and detach from GPU if applicable
    embeddings = embeddings.cpu().detach().numpy()
    
    # Normalize embeddings using FAISS
    faiss.normalize_L2(embeddings)
    
    end_time = time.time()
    # Uncomment the following line for timing information
    print(f"Embedding generation duration: {end_time - start_time} seconds")
    # Convert embeddings to a list of floats
    embeddings_list = embeddings[0].tolist()
    
    return embeddings_list



def get_openai_response(prompt):
    try:
        print('inside the response----------------')
        response_start_time = time.time()
        response = openai.ChatCompletion.create(
            # model="gpt-3.5-turbo",
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a knowledgeable subject matter expert."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            max_tokens=500,
            stop=None  # Adjust stop condition as needed
        )
        end_time = time.time()
        print(f"OpenAI API call duration: {end_time - response_start_time} --------------seconds")
        return response['choices'][0]['message']['content'].strip()
    except Exception as e:
        print(f"Error in OpenAI request: {e}")
        return "An error occurred while processing your request. Please try again."



def process_and_respond(request, user_question, document, conversation_messages):
    try:
        print('inside the process and response-------------------------')
        print(document,'----------------document')
        overall_start_time = time.time()

        # Retrieve chunks associated with the document
        chunks = document.chunks.all()

        # Extract embeddings from chunks
        chunk_embeddings = np.array([chunk.embeddings for chunk in chunks], dtype='float32')
        faiss.normalize_L2(chunk_embeddings)

        # Create FAISS index
        index = faiss.IndexFlatIP(chunk_embeddings.shape[1])
        index.add(chunk_embeddings)

        # Convert user question to embeddings using SentenceTransformer
        question_embedding = get_embeddings(user_question)
        question_embedding = np.array(question_embedding).reshape(1, -1).astype('float32')

        # Search for similar embeddings
        k = 5  # Number of similar chunks to retrieve
        distances, indices = index.search(question_embedding, k)
        similar_chunks_indices = indices[0].tolist()  # Convert int64 to regular list of integers

        # Filter out negative indices
        valid_indices = [int(i) for i in similar_chunks_indices if i >= 0]

        # Collect similar chunks' contents for context
        similar_chunks = [chunks[i] for i in valid_indices]
        context = " ".join([chunk.content for chunk in similar_chunks])

        # Append previous conversation messages to the context
        context += "\n\nPrevious Conversation:\n"
        for message in conversation_messages:
            context += f"{message['source'].capitalize()}: {message['text']}\n"

        # Generate the prompt using the template
        prompt = (
            "You are a knowledgeable medical expert with deep expertise in cardiology and other related domains of medical sciences. Use your expertise to answer user questions from the provided excerpts from the document."
            "If the information needed to answer the question is not in the excerpts, please indicate so. Provide a clear and concise answer. "
            "ensure the answers are in complete sentences"
            "Keep your response to fifteen sentences maximum.\n"
            f"Document Content: {context} \nQuestion: {user_question} \nAnswer:"
        )

        # Generate the response using OpenAI language model
        response = get_openai_response(prompt)

        # Calculate total response time
        overall_end_time = time.time()
        overall_duration_seconds = round(overall_end_time - overall_start_time, 2)
        overall_duration_display = f"{overall_duration_seconds} seconds" if overall_duration_seconds < 60 else f"{int(overall_duration_seconds // 60)} minutes and {round(overall_duration_seconds % 60, 2)} seconds"

        # Logging total response time
        print(f"Total process and respond duration: {overall_duration_display}")

        # Logging and session management
        request.session['conversation_messages'].append({'source': 'user', 'text': user_question})
        # request.session['conversation_messages'].append({'source': 'assistant', 'text': response})
        

        
        request.session.save()
        

        top_chunk_page_numbers = [chunk.page_number for chunk in similar_chunks]

        # Get unique page numbers while maintaining order
        seen = set()
        top_chunk_page_numbers = []
        for chunk in similar_chunks:
            if chunk.page_number not in seen:
                seen.add(chunk.page_number)
                top_chunk_page_numbers.append(chunk.page_number)

        print(top_chunk_page_numbers, '---------top_chunk_page_numbers')

        
        request.session['conversation_messages'].append({
            'source': 'assistant',
            'text': response,
            'page_numbers': top_chunk_page_numbers,
            'response_duration': overall_duration_display
        })
        request.session.save()
        print(request.session['conversation_messages'],'---------------------------update request.session')
        return response, overall_duration_display, top_chunk_page_numbers

    except Exception as e:
        print(f"Error in processing and responding: {e}")
        return "Error occurred while processing the question.", "0 seconds", []



from operator import attrgetter

def new_thread_id(doc_id):
    # get docuemnets with given id
    chat_conversation=ConversationChat.objects.filter(document_id=doc_id)
    if(len(chat_conversation) >0):
        # get max thread id for given docuemnt
        max_attr = max(chat_conversation, key=attrgetter('thread_id'))

        # increment the thread id and return
        return max_attr.thread_id + 1
    else:
        return 0
    pass


def delete_old_records_from_thread(docid, threadid):
    print('inside the delete_old records :doc id::--------------',docid)
    chat_conversation=ConversationChat.objects.filter(document_id=docid, thread_id = threadid)
    i = len(chat_conversation)
    max_attr = max(chat_conversation, key=attrgetter('id'))
    print(f' max_att: {type(max_attr.id)} --------------type---------------doc id : {type(docid)}')
    print(max_attr,'-------------------max_attr in doc chat ids')
    if max_attr is not None and i >1:
        print(' calling deleteions-----------------------')
        try:
            print('-------------- in try bolck')
            '''# print(chat_conversation.objects.all(),'--------------remaining---------records')
            # Delete all records except the one with the maximum ID
            #chat_conversation.objects.exclude(id=int(max_attr.chatid)).delete()
            # print(chat_conversation.objects.filter(id !=max_attr.chatid),'--------------filter---------records')
            #ConversationChat.objects.filter(id !=int(max_attr.id) , document_id=docid) .delete()
            # print(chat_conversation.objects.all(),'--------------remaining---------records')'''
            #ConversationChat.objects.filter(document_id=docid).exclude(id=max_attr.id, thread_id=max_attr.thread_id).delete()#
            
            ConversationChat.objects.filter(document_id=docid, thread_id = threadid).exclude(id=max_attr.id).delete()#
            # records_to_delete = ConversationChat.objects.exclude(id=max_attr.id, documentid=docid)
            # deleted_count, _ = records_to_delete.delete()

        except Exception as e:
            error_message = f"An error occurred: {e}"
            print('-------------- in except bolck')


def delete_old_records(docid):
    print('inside the delete_old records :doc id::--------------',docid)
    chat_conversation=ConversationChat.objects.filter(document_id=docid)
    i = len(chat_conversation)
    max_attr = max(chat_conversation, key=attrgetter('id'))
    print(f' max_att: {type(max_attr.id)} --------------type---------------doc id : {type(docid)}')
    print(max_attr,'-------------------max_attr in doc chat ids')
    if max_attr is not None and i >1:
        print(' calling deleteions-----------------------')
        try:
            print('-------------- in try bolck')
            '''# print(chat_conversation.objects.all(),'--------------remaining---------records')
            # Delete all records except the one with the maximum ID
            #chat_conversation.objects.exclude(id=int(max_attr.chatid)).delete()
            # print(chat_conversation.objects.filter(id !=max_attr.chatid),'--------------filter---------records')
            #ConversationChat.objects.filter(id !=int(max_attr.id) , document_id=docid) .delete()
            # print(chat_conversation.objects.all(),'--------------remaining---------records')'''
            #ConversationChat.objects.filter(document_id=docid).exclude(id=max_attr.id, thread_id=max_attr.thread_id).delete()#
            
            ConversationChat.objects.filter(document_id=docid).exclude(id=max_attr.id).delete()#
            # records_to_delete = ConversationChat.objects.exclude(id=max_attr.id, documentid=docid)
            # deleted_count, _ = records_to_delete.delete()

        except Exception as e:
            error_message = f"An error occurred: {e}"
            print('-------------- in except bolck')


from .forms import ConversationChatForm

@login_required
def chat_interface(request, document_id):
    OPENAI_API_KEY = settings.OPENAI_API_KEY

    try:
        document = get_object_or_404(UploadFile, pk=document_id)
        prompts = document.prompt.all()
        chat_conversation = ConversationChat.objects.filter(document_id=document_id)
        print(chat_conversation, '---------------------chat_conversation messages of document')

        chat_id = request.GET.get('chat_id')
        tid = request.GET.get('thread_id')
        specific_chat = None

        if chat_id:
            specific_chat = get_object_or_404(ConversationChat, id=chat_id)
            request.session['to_delete_chat_id'] = chat_id
            request.session['conversation_messages'] = specific_chat.messages
            print()
            print(chat_id,'------------------------chat_id history ')
            print(request.session['conversation_messages'],'--------------conversation_messages after click chat history--------')

        if tid:
            thread_id = tid
        else:
            thread_id = -1

        request.session['nti'] = thread_id

        if 'document_id' in request.session and request.session['document_id'] != document_id:
            request.session['conversation_messages'] = []
            request.session['last_messages'] = []
            request.session['nti'] = -1

            if len(chat_conversation) == 0:
                thread_id = 0
            else:
                thread_id = new_thread_id(document_id)

            request.session['nti'] = thread_id
            request.session.save()

        request.session['document_id'] = document_id
        request.session.save()

        if 'conversation_messages' not in request.session:
            request.session['conversation_messages'] = []

        if 'last_messages' not in request.session:
            request.session['last_messages'] = []

        share_url = request.build_absolute_uri()

        if request.method == 'POST':
            if 'edit_title' in request.POST:
                chat_id = request.POST.get('chat_id')
                chat = get_object_or_404(ConversationChat, id=chat_id)
                form = ConversationChatForm(request.POST, instance=chat)
                if form.is_valid():
                    form.save()
                    return redirect('chat_interface', document_id=document.id)

            if 'submit_action' in request.POST:
                submit_action = request.POST.get('submit_action')

                if submit_action == 'prompt':
                    prompt_data = request.POST.get('prompt_data')
                    user_question = prompt_data
                    if len(user_question) >= 5:
                        # Update the context with the current conversation history
                        new_response = process_and_respond(request, user_question, document, request.session.get('conversation_messages', []))
                        request.session['conversation_messages'] = request.session.get('conversation_messages', []) + [{'source': 'user', 'text': user_question}, {'source': 'assistant', 'text': new_response}]
                        request.session.save()
                        return JsonResponse({'response_prompt': new_response})

                elif submit_action == 'submit':
                    user_question = request.POST.get('question')
                    print(user_question, '----------------')
                    thread_id = request.POST.get('thread_id')
                    request.session['nti'] = thread_id
                    request.session.save()
                    if len(user_question) >= 5:
                        last_messages = request.session.get('conversation_messages', [])
                        response_text, response_time, page_num = process_and_respond(request, user_question, document, last_messages)
                        request.session['nti'] = thread_id
                        request.session.save()
                        return JsonResponse({'response_prompt': response_text, 'page_num': page_num, 'response_time': response_time})
                    else:
                        messages.error(request, "User input must contain a minimum of 6 characters.")
                        return render(request, 'user/chat_interface.html', {
                            'conversation_messages': request.session.get('conversation_messages', []),
                            'document': document,
                            'prompts': prompts,
                            'chat_conversation': chat_conversation,
                            'thread_id': thread_id,
                            'specific_chat': specific_chat,
                            'form': ConversationChatForm(instance=specific_chat) if specific_chat else None
                        })

                elif submit_action == 'clear':
                    request.session['conversation_messages'] = []
                    request.session.save()

            return render(request, 'user/chat_interface.html', {
                'conversation_messages': request.session.get('conversation_messages', []),
                'document': document,
                'prompts': prompts,
                'share_url': share_url,
                'chat_conversation': chat_conversation,
                'thread_id': thread_id,
                'specific_chat': specific_chat,
                'form': ConversationChatForm(instance=specific_chat) if specific_chat else None
            })

        return render(request, 'user/chat_interface.html', {
            'conversation_messages': request.session['conversation_messages'],
            'document': document,
            'prompts': prompts,
            'share_url': share_url,
            'chat_conversation': chat_conversation,
            'thread_id': thread_id,
            'specific_chat': specific_chat,
            'form': ConversationChatForm(instance=specific_chat) if specific_chat else None
        })

    except Exception as e:
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message, 'chat_conversation': chat_conversation})


from django.http import JsonResponse, HttpResponseRedirect
from django.urls import reverse

from django.views.decorators.http import require_POST
from .models import ConversationChat






@login_required
@require_POST
def save_chat(request):
    if request.method == 'POST':
        try:
            # Get conversation messages from the session
            conversation_messages = request.session.get('conversation_messages', [])
            chat_id = request.session.get('chat_id')
            print(chat_id,'-------------chat -id-')
            print(conversation_messages, '------------------------request.session[conversation_messages]')

            # Retrieve the document associated with the provided ID
            document_id = request.session.get('document_id')
            nti = request.session.get('nti')
            # print("In save chat,, received thead id", nti)
            if not document_id:
                return JsonResponse({'success': False, 'message': 'Document ID is missing.'})

            document = UploadFile.objects.get(pk=document_id)
            response_duration = ''
            page_numbers = []
            src_name = ''
            if conversation_messages:
                last_message = conversation_messages[-1]
                if last_message['source'] == 'assistant':
                    response_duration = last_message.get('response_duration', '')
                    page_numbers = last_message.get('page_numbers', [])
                    # src_name = last_message.get('source', '')

            print(page_numbers,'-page_numbers  in save chat------------')
            print(response_duration,'-response_duration  in save chat------------')
            # Save chat history using the ConversationChat model
            new_chat = ConversationChat.objects.create(
                user=request.user,
                document=document,
                messages=conversation_messages,
                thread_id = nti
                
            )

           
            del_rec=delete_old_records_from_thread(document_id, nti)
            print(del_rec,'------------------deleted old_ones returns max id inside the save chat view')

            # Redirect to the user dashboard
            return HttpResponseRedirect(reverse('user_dashboard'))

        except Exception as e:
            return JsonResponse({'success': False, 'message': str(e)})





@login_required
def common_chat_interface(request):
    try:
        # Reset session data for a new session
        request.session['common_conversation_messages'] = []
        request.session['last_messages'] = []
        
        OPENAI_API_KEY = settings.OPENAI_API_KEY
        documents = UploadFile.objects.all()
        prompts = Prompt.objects.all()
        
        # Get the absolute URL for sharing
        share_url = request.build_absolute_uri()

        if request.method == 'POST':
            if 'submit_action' in request.POST:
                submit_action = request.POST.get('submit_action')

                if submit_action == 'prompt':
                    user_question = request.POST.get('prompt_data')
                elif submit_action == 'submit':
                    user_question = request.POST.get('question')

                if user_question and len(user_question) >= 5:
                    new_responses = []

                    # Use ThreadPoolExecutor to run process_and_respond concurrently
                    with ThreadPoolExecutor() as executor:
                        futures = [
                            executor.submit(common_process_and_respond, request, user_question, doc, request.session.get('last_messages', []))
                            for doc in documents
                        ]

                        for future in futures:
                            try:
                                response_text, doc_title, duration, page_numbers = future.result()
                                new_responses.append((response_text, doc_title, duration, page_numbers))
                            except Exception as e:
                                print(f"Error processing document: {e}")

                    # Extract response texts, document titles, durations, and page numbers
                    response_texts = [response[0] for response in new_responses]
                    document_titles = [response[1] for response in new_responses]
                    durations = [response[2] for response in new_responses]
                    page_numbers_list = [response[3] for response in new_responses]

                    print(document_titles, '----------------document_titles')
                    print(durations, '----------------durations')

                    # Update session messages
                    for doc_response in new_responses:
                        response_text, doc_title, duration, page_numbers = doc_response
                        request.session['common_conversation_messages'].append({
                            'document': doc_title,
                            'source': 'assistant',
                            'text': response_text,
                            'duration': duration,
                            'page_numbers': page_numbers
                        })
                    request.session.save()

                    return JsonResponse({
                        'response_texts': response_texts,
                        'document_titles': document_titles,
                        'durations': durations,
                        'page_numbers_list': page_numbers_list
                    })

                elif submit_action == 'clear':
                    # Clear conversation messages stored in session
                    request.session['common_conversation_messages'] = []
                    request.session.save()
                else:
                    messages.error(request, "User input must contain a minimum of 6 characters.")
                    return render(request, 'user/common_chat_interface.html', {
                        'common_conversation_messages': request.session.get('common_conversation_messages', []),
                        'documents': documents,
                        'prompts': prompts
                    })

        return render(request, 'user/common_chat_interface.html', {
            'common_conversation_messages': request.session.get('common_conversation_messages', []),
            'documents': documents,
            'prompts': prompts,
            'share_url': share_url
        })

    except Exception as e:
        error_message = f"An error occurred: {e}"
        return render(request, 'user/common_chat_interface.html', {'error_message': error_message})


def common_process_and_respond(request, user_question, document, last_messages):
    try:
        print('inside the process and response-------------------------')
        overall_start_time = time.time()

        # Retrieve chunks associated with the document
        chunks = document.chunks.all()

        # Extract embeddings from chunks
        chunk_embeddings = np.array([chunk.embeddings for chunk in chunks], dtype='float32')
        faiss.normalize_L2(chunk_embeddings)

        # Create FAISS index
        index = faiss.IndexFlatIP(chunk_embeddings.shape[1])
        index.add(chunk_embeddings)

        # Convert user question to embeddings using SentenceTransformer
        question_embedding = get_embeddings(user_question)
        question_embedding = np.array(question_embedding).reshape(1, -1).astype('float32')

        # Search for similar embeddings
        k = 5  # Number of similar chunks to retrieve
        distances, indices = index.search(question_embedding, k)
        similar_chunks_indices = indices[0].tolist()  # Convert int64 to regular list of integers

        # Filter out negative indices
        valid_indices = [int(i) for i in similar_chunks_indices if i >= 0]

        # Collect similar chunks' contents for context
        similar_chunks = [chunks[i] for i in valid_indices]
        context = " ".join([chunk.content for chunk in similar_chunks])

        # Generate the prompt using the template
        prompt = (
            "You are a knowledgeable medical expert with deep expertise in cardiology and other related domains of medical sciences. Use your expertise to answer user questions from the provided excerpts from the document."
            "If the information needed to answer the question is not in the excerpts or not in the document  please indicate so. Provide a clear and concise answer. "
            "ensure the answers are in complete sentences"
            "Ignore repetitive text patterns "
            "Keep your response to fifteen sentences maximum.\n"
            f"Document Content: {context} \nQuestion: {user_question} \nAnswer:"
        )

        # Generate the response using OpenAI language model
        response = get_openai_response(prompt)

        # Calculate total response time
        overall_end_time = time.time()
        overall_duration_seconds = round(overall_end_time - overall_start_time, 2)
        overall_duration_display = f"{overall_duration_seconds} seconds" if overall_duration_seconds < 60 else f"{int(overall_duration_seconds // 60)} minutes and {round(overall_duration_seconds % 60, 2)} seconds"

        # Logging total response time
        print(f"Total process and respond duration: {overall_duration_display}")

        # Logging and session management
        request.session['common_conversation_messages'].append({'source': 'user', 'text': user_question})
        request.session['common_conversation_messages'].append({'source': 'assistant', 'text': response})
        request.session.save()

        # Get the document title
        doc_title = document.title

        top_chunk_page_numbers = [chunk.page_number for chunk in similar_chunks]

        return response, doc_title, overall_duration_display, top_chunk_page_numbers

    except Exception as e:
        print(f"Error in processing and responding: {e}")
        return "Error occurred while processing the question.", "Unknown Document", "0 seconds", []




def filter_response(response, chunks):
    # Calculate relevance score for each response
    relevance_scores = []
    print('-----------------inside the filter_response')
    for chunk in chunks:
        relevance_score = calculate_relevance_score(response, chunk.content)
        relevance_scores.append(relevance_score)

    # Find the chunk with the highest relevance score
    max_relevance_score_index = relevance_scores.index(max(relevance_scores))
    max_relevance_chunk = chunks[max_relevance_score_index]
    page_num = max_relevance_chunk.page_number

    # Threshold filtering
    threshold = 0.9  # Adjust threshold as needed
    if max(relevance_scores) >= threshold:
        return page_num
    else:
        return page_num



def calculate_relevance_score(response, chunk_content):
    # Calculate relevance score based on how much of the chunk content is present in the response
    overlap = len(set(response.split()) & set(chunk_content.split()))
    relevance_score = overlap / len(chunk_content.split())  # Normalize by chunk length
    return relevance_score



@login_required
def copy_to_clipboard(request):
    if request.method == 'GET':
        # Retrieve the AI response text from the request
        ai_response_text = request.GET.get('ai_response_text')
        
        if ai_response_text:
            try:
                # Use pyperclip to copy the text to the clipboard
                pyperclip.copy(ai_response_text)
                return JsonResponse({'success': True, 'message': 'Text copied to clipboard successfully'})
            except Exception as e:
                return JsonResponse({'success': False, 'error': f'Failed to copy text to clipboard: {str(e)}'}, status=500)
        
        # If no text is copied or the method is not GET, return an error response
        return JsonResponse({'success': False, 'error': 'No text provided'}, status=400)



def clear_chat(request):
    # Clear conversation messages stored in session
    print('inside the clear fun')
    request.session['conversation_messages'] = []
    request.session.save()
    return JsonResponse({'success': True, 'message': 'Chat cleared successfully'})




def document_viewer(request, document_id):
    document = get_object_or_404(UploadFile, id=document_id)
    return render(request, 'user/document_viewer.html', {'document': document})




#for normal conversation

def wrap_text(text, max_width, pdf_canvas):
    """Wrap text to fit within the given width."""
    lines = []
    words = text.split(' ')
    current_line = words[0]
    for word in words[1:]:
        test_line = f"{current_line} {word}"
        width = pdf_canvas.stringWidth(test_line, 'Helvetica', 12)
        if width <= max_width:
            current_line = test_line
        else:
            lines.append(current_line)
            current_line = word
    lines.append(current_line)
    return lines


from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader

def add_watermark_pdf(pdf_canvas, width, height, text=None, image_path=None):
    if text:
        pdf_canvas.saveState()
        pdf_canvas.setFont("Helvetica", 40)
        pdf_canvas.setFillGray(0.5, 0.5)
        pdf_canvas.translate(width / 2, height / 2)
        pdf_canvas.rotate(45)
        pdf_canvas.drawCentredString(0, 0, text)
        pdf_canvas.restoreState()

    if image_path and os.path.exists(image_path):
        pdf_canvas.saveState()
        pdf_canvas.setFillAlpha(0.1)
        pdf_canvas.drawImage(image_path, width / 4, height / 4, width / 2, height / 2, mask='auto')
        pdf_canvas.restoreState()

@login_required
def export_conversation_pdf(request):
    try:
        conversation_messages = request.session.get('conversation_messages', [])
        
        if not conversation_messages:
            raise ValueError("No conversation messages found in session.")

        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.pdf'

        pdf_canvas = canvas.Canvas(response, pagesize=letter)
        width, height = letter

        # Add watermark (text or image)
        watermark_text = " "
        # watermark_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'logos', 'centrix_fav.png')  # Adjust the path accordingly
        watermark_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'logos', 'Medask-Watermark.jpg')
        add_watermark_pdf(pdf_canvas, width, height, text=watermark_text, image_path=watermark_image_path)

        pdf_canvas.setFont("Helvetica-Bold", 16)
        pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
        pdf_canvas.drawString(100, height - 40, "Conversation Export")

        y_position = height - 60
        line_height = 14
        max_width = width - 80

        for message in conversation_messages:
            if message['source'] == 'user':
                pdf_canvas.setFont("Helvetica-Bold", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                pdf_canvas.drawString(40, y_position, "User:")
                y_position -= line_height
                pdf_canvas.setFont("Helvetica", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                lines = wrap_text(message['text'], max_width, pdf_canvas)
                for line in lines:
                    pdf_canvas.drawString(40, y_position, line)
                    y_position -= line_height
            elif message['source'] == 'assistant':
                pdf_canvas.setFont("Helvetica-Bold", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                pdf_canvas.drawString(40, y_position, "Assistant:")
                y_position -= line_height
                pdf_canvas.setFont("Helvetica", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                lines = wrap_text(message['text'], max_width, pdf_canvas)
                for line in lines:
                    pdf_canvas.drawString(40, y_position, line)
                    y_position -= line_height

            y_position -= line_height

            if y_position < 40:
                pdf_canvas.showPage()
                add_watermark_pdf(pdf_canvas, width, height, text=watermark_text, image_path=watermark_image_path)  # Add watermark on new page
                y_position = height - 40

        pdf_canvas.save()
        return response

    except Exception as e:
        logging.error("Error exporting conversation to PDF: %s", e)
        return HttpResponse(status=500)




  

from docx import Document

@login_required
def export_conversation_doc(request):
    try:
        conversation_messages = request.session.get('conversation_messages', [])
        print("Conversation messages in doc export: ", conversation_messages)
        
        if not conversation_messages:
            raise ValueError("No conversation messages found in session.")

        doc = Document()
        doc.add_heading('Conversation Export', 0)

        for message in conversation_messages:
            if message['source'] == 'user':
                doc.add_heading('User:', level=1)
                doc.add_paragraph(message['text'])
            elif message['source'] == 'assistant':
                doc.add_heading('Assistant:', level=1)
                doc.add_paragraph(message['text'])
                

        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.docx'
        doc.save(response)

        return response

    except Exception as e:
        logging.error(f"Error exporting conversation: {e}")
        return render(request, 'user/chat_interface.html', {'error_message': str(e)})




@login_required
def export_conversation_json(request):
    try:
        # Retrieve conversation messages from session
        conversation_messages = request.session.get('conversation_messages', [])

        if not conversation_messages:
            raise ValueError("No conversation messages found in session.")

        # Create a JSON response
        response = HttpResponse(content_type='application/json')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.json'

        # Write conversation messages to JSON
        response.write(json.dumps(conversation_messages, indent=4))

        return response

    except Exception as e:
        # Log the error for debugging
        logging.error(f"Error exporting conversation: {e}")

        # Render the chat interface with an error message
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})





@login_required
def export_conversation_ppt(request):
    try:
        # Retrieve conversation messages from session
        conversation_messages = request.session.get('conversation_messages', [])

        if not conversation_messages:
            raise ValueError("No conversation messages found in session.")

        # Create a new Presentation
        prs = Presentation()

        # Define the background image path
        # background_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'ppt_template_img', 'pptbg.jpg')
        background_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'ppt_template_img', 'Medask_PPT_theme.jpg')

        if not os.path.exists(background_image_path):
            raise ValueError(f"Background image not found at {background_image_path}")

        # Get slide dimensions from the presentation object
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Function to set background image for a slide
        def set_background_image(slide, image_path, slide_width, slide_height):
            left = top = 0
            pic = slide.shapes.add_picture(image_path, left, top, width=slide_width, height=slide_height)
            slide.shapes._spTree.remove(pic._element)  # Remove the added picture element
            slide.shapes._spTree.insert(2, pic._element)  # Insert the picture as the background

        # Function to split text into slides with a maximum of max_words_per_slide words per slide
        def split_text_by_word_limit(text, max_words_per_slide):
            words = text.split()
            slides_content = []
            current_slide_text = []
            current_word_count = 0

            for word in words:
                current_slide_text.append(word)
                current_word_count += 1

                if current_word_count >= max_words_per_slide:
                    # Ensure the slide ends with a meaningful completion
                    last_char = word[-1]
                    if last_char in '.!?':
                        slides_content.append(' '.join(current_slide_text).strip())
                        current_slide_text = []
                        current_word_count = 0

            if current_slide_text:
                slides_content.append(' '.join(current_slide_text).strip())

            return slides_content

        # Add title slide with background image
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        set_background_image(slide, background_image_path, slide_width, slide_height)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Conversation Export"
        subtitle.text = "Exported conversation messages"

        # Main Title Style
        title.text_frame.paragraphs[0].font.size = Pt(38)
        title.text_frame.paragraphs[0].font.name = 'Calibri'
        title.text_frame.paragraphs[0].font.bold = True

        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
        subtitle.text_frame.paragraphs[0].font.bold = True

        max_words_per_slide = 80
        last_user_message = None

        # Add conversation messages to slides with background image
        for message in conversation_messages:
            if message['source'] == 'user':
                last_user_message = message['text']
            else:  # message['source'] == 'assistant'
                if last_user_message:
                    title_text = last_user_message
                    last_user_message = None
                else:
                    title_text = "Assistant"

                content = f"{message['text']}"
                slides_content = split_text_by_word_limit(content, max_words_per_slide)
                for slide_text in slides_content:
                    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
                    slide = prs.slides.add_slide(slide_layout)
                    set_background_image(slide, background_image_path, slide_width, slide_height)
                    title = slide.shapes.title
                    title.text = title_text

                    # Slide Title Style
                    title.text_frame.paragraphs[0].font.size = Pt(28)
                    title.text_frame.paragraphs[0].font.name = 'Calibri'
                    title.text_frame.paragraphs[0].font.bold = True

                    content_box = slide.shapes.placeholders[1]
                    content_box.text = slide_text

                    # Set font size, style, and type for the content
                    for paragraph in content_box.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(20)
                            run.font.name = 'Calibri'
                            run.font.bold = False

                    # Add watermark text
                    watermark_text = "Centrix Technology copyright"
                    watermark = slide.shapes.add_textbox(left=Pt(10), top=Pt(10), width=Pt(500), height=Pt(50))
                    watermark.text_frame.text = watermark_text
                    for paragraph in watermark.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(20)
                            run.font.color.rgb = RGBColor(192, 192, 192)  # Light gray color

        # Add closing slide with concluding message
        closing_slide_layout = prs.slide_layouts[1]  # Use the title and content layout for the closing slide
        slide = prs.slides.add_slide(closing_slide_layout)
        set_background_image(slide, background_image_path, slide_width, slide_height)
        title = slide.shapes.title
        title.text = "Conclusion"

        # Closing Title Style
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = 'Calibri'
        title.text_frame.paragraphs[0].font.bold = True

        content_box = slide.shapes.placeholders[1]
        content_box.text = "This concludes the export of conversation messages. Thank you for viewing."

        # Set font size, style, and type for the closing slide content
        for paragraph in content_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.name = 'Calibri'
                run.font.bold = False

        # Create a response with the presentation as an attachment
        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.pptx'
        prs.save(response)

        return response

    except Exception as e:
        # Log the error for debugging
        logging.error(f"Error exporting conversation: {e}")

        # Render the chat interface with an error message
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})


# export to audio 

# language:english (US)
#voice: joanna
# gender :Female


import boto3
from botocore.exceptions import BotoCoreError, ClientError
from django.shortcuts import render
from django.http import HttpResponse
from django.contrib.auth.decorators import login_required
from django.utils.text import slugify
from contextlib import closing
import sys

@login_required
def export_conversation_audio(request):
    try:
        conversation_messages = request.session.get('conversation_messages', [])
        print(conversation_messages,'------------------conversation_messages')
        
        if not conversation_messages:
            raise ValueError("No conversation messages found in session.")
        
        print(f"Conversation messages: {conversation_messages}")
        sys.stdout.flush()
        aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
        # Initialize Polly client
        polly_client = boto3.Session(
            aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            region_name='ap-southeast-1'
        ).client('polly')
        print(aws_access_key_id,'-------------------------aws_access_key_id')
        print("Initialized AWS Polly client.")
        sys.stdout.flush()

        # Combine all messages into a single text
        full_text = ""
        for message in conversation_messages:
            source = "User" if message['source'] == 'user' else "Assistant"
            full_text += f"{source}: {message['text']}\n"
        
        print(f"Full text for synthesis: {full_text}")
        sys.stdout.flush()

        # Call Polly to synthesize the text
        response = polly_client.synthesize_speech(
            Text=full_text,
            OutputFormat='mp3',
            VoiceId='Kajal',  # Using Joanna voice
            LanguageCode='en-IN',
            Engine='neural'
        )
  

        print("Called AWS Polly to synthesize speech.")
        sys.stdout.flush()

        if "AudioStream" in response:
            with closing(response["AudioStream"]) as stream:
                audio_content = stream.read()
                response = HttpResponse(content_type='audio/mpeg')
                response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.mp3'
                response.write(audio_content)
                print("Successfully created audio response.")
                sys.stdout.flush()
                # Print the size of the audio content
                print(f"Size of audio content: {len(audio_content)} bytes")
                sys.stdout.flush()
                return response
        else:
            raise ValueError("Could not synthesize speech using AWS Polly.")

    except (BotoCoreError, ClientError) as error:
        print(f"Error calling AWS Polly: {error}")
        sys.stdout.flush()
        error_message = f"An error occurred while synthesizing audio: {error}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})
    except Exception as e:
        print(f"Error exporting conversation to audio: {e}")
        sys.stdout.flush()
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})



#for common chat conversation export_conversation




# for common chat export [comm conversation messages]




@login_required
def common_export_conversation_pdf(request):
    try:
        common_conversation_messages = request.session.get('common_conversation_messages', [])
        
        if not common_conversation_messages:
            raise ValueError("No conversation messages found in session.")

        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.pdf'

        pdf_canvas = canvas.Canvas(response, pagesize=letter)
        width, height = letter

        # Add watermark (text or image)
        watermark_text = " "# add text here
        watermark_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'logos', 'Medask-Watermark.jpg')  # Adjust the path accordingly
        add_watermark_pdf(pdf_canvas, width, height, text=watermark_text, image_path=watermark_image_path)

        pdf_canvas.setFont("Helvetica-Bold", 16)
        pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
        pdf_canvas.drawString(100, height - 40, "Conversation Export")

        y_position = height - 60
        line_height = 14
        max_width = width - 80

        for message in common_conversation_messages:
            if message['source'] == 'user':
                pdf_canvas.setFont("Helvetica-Bold", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                pdf_canvas.drawString(40, y_position, "User:")
                y_position -= line_height
                pdf_canvas.setFont("Helvetica", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                lines = wrap_text(message['text'], max_width, pdf_canvas)
                for line in lines:
                    pdf_canvas.drawString(40, y_position, line)
                    y_position -= line_height
            elif message['source'] == 'assistant':
                pdf_canvas.setFont("Helvetica-Bold", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                pdf_canvas.drawString(40, y_position, "Assistant:")
                y_position -= line_height
                pdf_canvas.setFont("Helvetica", 12)
                pdf_canvas.setFillColorRGB(0, 0, 0)  # Solid black color
                lines = wrap_text(message['text'], max_width, pdf_canvas)
                for line in lines:
                    pdf_canvas.drawString(40, y_position, line)
                    y_position -= line_height

            y_position -= line_height

            if y_position < 40:
                pdf_canvas.showPage()
                add_watermark_pdf(pdf_canvas, width, height, text=watermark_text, image_path=watermark_image_path)  # Add watermark on new page
                y_position = height - 40

        pdf_canvas.save()
        return response

    except Exception as e:
        logging.error("Error exporting conversation to PDF: %s", e)
        return HttpResponse(status=500)




  

@login_required
def common_export_conversation_doc(request):
    try:
        common_conversation_messages = request.session.get('common_conversation_messages', [])
        print("Conversation messages in doc export: ", common_conversation_messages)
        
        if not common_conversation_messages:
            raise ValueError("No conversation messages found in session.")

        doc = Document()
        doc.add_heading('Conversation Export', 0)

        for message in common_conversation_messages:
            if message['source'] == 'user':
                doc.add_heading('User:', level=1)
                doc.add_paragraph(message['text'])
            elif message['source'] == 'assistant':
                doc.add_heading('Assistant:', level=1)
                doc.add_paragraph(message['text'])

        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.docx'
        doc.save(response)

        return response

    except Exception as e:
        logging.error(f"Error exporting conversation: {e}")
        return render(request, 'user/chat_interface.html', {'error_message': str(e)})



@login_required
def common_export_conversation_json(request):
    try:
        # Retrieve conversation messages from session
        common_conversation_messages = request.session.get('common_conversation_messages', [])

        if not common_conversation_messages:
            raise ValueError("No conversation messages found in session.")

        # Create a JSON response
        response = HttpResponse(content_type='application/json')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.json'

        # Write conversation messages to JSON
        response.write(json.dumps(common_conversation_messages, indent=4))

        return response

    except Exception as e:
        # Log the error for debugging
        logging.error(f"Error exporting conversation: {e}")

        # Render the chat interface with an error message
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})





@login_required
def common_export_conversation_ppt(request):
    try:
        # Retrieve conversation messages from session
        common_conversation_messages = request.session.get('common_conversation_messages', [])

        if not common_conversation_messages:
            raise ValueError("No conversation messages found in session.")

        # Create a new Presentation
        prs = Presentation()

        # Define the background image path
        # background_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'ppt_template_img', 'pptbg.jpg')
        background_image_path = os.path.join(settings.BASE_DIR, 'static', 'media', 'ppt_template_img', 'Medask_PPT_theme.jpg')

        if not os.path.exists(background_image_path):
            raise ValueError(f"Background image not found at {background_image_path}")

        # Get slide dimensions from the presentation object
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Function to set background image for a slide
        def set_background_image(slide, image_path, slide_width, slide_height):
            left = top = 0
            pic = slide.shapes.add_picture(image_path, left, top, width=slide_width, height=slide_height)
            slide.shapes._spTree.remove(pic._element)  # Remove the added picture element
            slide.shapes._spTree.insert(2, pic._element)  # Insert the picture as the background

        # Function to split text into slides with a maximum of max_words_per_slide words per slide
        def split_text_by_word_limit(text, max_words_per_slide):
            words = text.split()
            slides_content = []
            current_slide_text = []
            current_word_count = 0

            for word in words:
                current_slide_text.append(word)
                current_word_count += 1

                if current_word_count >= max_words_per_slide:
                    # Ensure the slide ends with a meaningful completion
                    last_char = word[-1]
                    if last_char in '.!?':
                        slides_content.append(' '.join(current_slide_text).strip())
                        current_slide_text = []
                        current_word_count = 0

            if current_slide_text:
                slides_content.append(' '.join(current_slide_text).strip())

            return slides_content

        # Add title slide with background image
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        set_background_image(slide, background_image_path, slide_width, slide_height)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Conversation Export"
        subtitle.text = "Exported conversation messages"

        # Main Title Style
        title.text_frame.paragraphs[0].font.size = Pt(38)
        title.text_frame.paragraphs[0].font.name = 'Calibri'
        title.text_frame.paragraphs[0].font.bold = True

        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
        subtitle.text_frame.paragraphs[0].font.bold = True

        max_words_per_slide = 80
        last_user_message = None

        # Add conversation messages to slides with background image
        for message in common_conversation_messages:
            if message['source'] == 'user':
                last_user_message = message['text']
            else:  # message['source'] == 'assistant'
                if last_user_message:
                    title_text = last_user_message
                    last_user_message = None
                else:
                    title_text = "Assistant"

                content = f"{message['text']}"
                slides_content = split_text_by_word_limit(content, max_words_per_slide)
                for slide_text in slides_content:
                    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
                    slide = prs.slides.add_slide(slide_layout)
                    set_background_image(slide, background_image_path, slide_width, slide_height)
                    title = slide.shapes.title
                    title.text = title_text

                    # Slide Title Style
                    title.text_frame.paragraphs[0].font.size = Pt(28)
                    title.text_frame.paragraphs[0].font.name = 'Calibri'
                    title.text_frame.paragraphs[0].font.bold = True

                    content_box = slide.shapes.placeholders[1]
                    content_box.text = slide_text

                    # Set font size, style, and type for the content
                    for paragraph in content_box.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(20)
                            run.font.name = 'Calibri'
                            run.font.bold = False

                    # Add watermark text
                    watermark_text = "Centrix Technology copyright"
                    watermark = slide.shapes.add_textbox(left=Pt(10), top=Pt(10), width=Pt(500), height=Pt(50))
                    watermark.text_frame.text = watermark_text
                    for paragraph in watermark.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(20)
                            run.font.color.rgb = RGBColor(192, 192, 192)  # Light gray color

        # Add closing slide with concluding message
        closing_slide_layout = prs.slide_layouts[1]  # Use the title and content layout for the closing slide
        slide = prs.slides.add_slide(closing_slide_layout)
        set_background_image(slide, background_image_path, slide_width, slide_height)
        title = slide.shapes.title
        title.text = "Conclusion"

        # Closing Title Style
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.name = 'Calibri'
        title.text_frame.paragraphs[0].font.bold = True

        content_box = slide.shapes.placeholders[1]
        content_box.text = "This concludes the export of conversation messages. Thank you for viewing."

        # Set font size, style, and type for the closing slide content
        for paragraph in content_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.name = 'Calibri'
                run.font.bold = False

        # Create a response with the presentation as an attachment
        response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.pptx'
        prs.save(response)

        return response

    except Exception as e:
        # Log the error for debugging
        logging.error(f"Error exporting conversation: {e}")

        # Render the chat interface with an error message
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})



@login_required
def common_export_conversation_audio(request):
    try:
        common_conversation_messages = request.session.get('common_conversation_messages', [])
        print(common_conversation_messages,'------------------common_conversation_messages')
        
        if not common_conversation_messages:
            raise ValueError("No conversation messages found in session.")
        
        print(f"Conversation messages: {common_conversation_messages}")
        sys.stdout.flush()
        aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
        # Initialize Polly client
        polly_client = boto3.Session(
            aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
            aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
            region_name='ap-southeast-1'
        ).client('polly')
        print(aws_access_key_id,'-------------------------aws_access_key_id')
        print("Initialized AWS Polly client.")
        sys.stdout.flush()

        # Combine all messages into a single text
        full_text = ""
        for message in common_conversation_messages:
            source = "User" if message['source'] == 'user' else "Assistant"
            full_text += f"{source}: {message['text']}\n"
        
        print(f"Full text for synthesis: {full_text}")
        sys.stdout.flush()

        # Call Polly to synthesize the text
        response = polly_client.synthesize_speech(
            Text=full_text,
            OutputFormat='mp3',
            VoiceId='Kajal',  # Using Joanna voice
            LanguageCode='en-IN',
            Engine='neural'
        )
  

        print("Called AWS Polly to synthesize speech.")
        sys.stdout.flush()

        if "AudioStream" in response:
            with closing(response["AudioStream"]) as stream:
                audio_content = stream.read()
                response = HttpResponse(content_type='audio/mpeg')
                response['Content-Disposition'] = f'attachment; filename={slugify("conversation_export")}.mp3'
                response.write(audio_content)
                print("Successfully created audio response.")
                sys.stdout.flush()
                # Print the size of the audio content
                print(f"Size of audio content: {len(audio_content)} bytes")
                sys.stdout.flush()
                return response
        else:
            raise ValueError("Could not synthesize speech using AWS Polly.")

    except (BotoCoreError, ClientError) as error:
        print(f"Error calling AWS Polly: {error}")
        sys.stdout.flush()
        error_message = f"An error occurred while synthesizing audio: {error}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})
    except Exception as e:
        print(f"Error exporting conversation to audio: {e}")
        sys.stdout.flush()
        error_message = f"An error occurred: {e}"
        return render(request, 'user/chat_interface.html', {'error_message': error_message})
